import asyncio
import json
import os
import math
import random
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional

from config import (
    AI_PROVIDER,
    GEMINI_API_KEY,
    GEMINI_MODEL,
    OPENAI_API_KEY,
    OPENAI_MODEL,
)


class AIServiceError(RuntimeError):
    pass


def _looks_like_openai_key(key: Optional[str]) -> bool:
    k = (key or "").strip()
    return k.startswith("sk-")


def _looks_like_gemini_key(key: Optional[str]) -> bool:
    k = (key or "").strip()
    return k.startswith("AIza")


def _looks_like_gemini_leaked_key_error(err_text: str) -> bool:
    low = (err_text or "").lower()
    return "reported as leaked" in low or ("leaked" in low and "api key" in low)


def _looks_like_gemini_auth_error(err_text: str) -> bool:
    low = (err_text or "").lower()
    return (
        "api key" in low
        or "403" in low
        or "401" in low
        or "permission" in low
        or "unauthorized" in low
        or "forbidden" in low
    )


def _format_gemini_auth_error(err_text: str) -> str:
    if _looks_like_gemini_leaked_key_error(err_text):
        return (
            "Gemini 403: bu API key Google tomonidan leaked deb belgilangan va bloklangan.\n"
            "Yechim: yangi Gemini API key yarating, `.env` dagi `GEMINI_API_KEY` ni almashtiring va botni qayta ishga tushiring.\n"
            "Agar OpenAI ishlatmasangiz, `.env` dagi `OPENAI_API_KEY` ni olib tashlang."
        )
    return "Gemini API key noto'g'ri yoki ruxsat yo'q. `.env` dagi `GEMINI_API_KEY` ni tekshiring."


def _gemini_retry_after_seconds(err_text: str) -> Optional[int]:
    low = (err_text or "").lower()
    # Example: "Please retry in 41.8055s."
    m = re.search(r"retry\s+in\s+([0-9]+(?:\.[0-9]+)?)s", low)
    if m:
        try:
            return max(0, int(round(float(m.group(1)))))
        except Exception:
            pass
    # Example: "retry_delay { seconds: 41 }"
    m = re.search(r"retry_delay\s*\{[^}]*seconds\s*:\s*([0-9]+)\s*\}", low, flags=re.DOTALL)
    if m and m.group(1).isdigit():
        return max(0, int(m.group(1)))
    return None


def _gemini_quota_value(err_text: str) -> Optional[int]:
    m = re.search(r"quota_value\s*:\s*([0-9]+)", err_text or "", flags=re.IGNORECASE)
    if m and m.group(1).isdigit():
        try:
            return int(m.group(1))
        except Exception:
            return None
    return None


def _looks_like_gemini_quota_error(err_text: str) -> bool:
    low = (err_text or "").lower()
    return (
        " 429" in low
        or low.startswith("429")
        or "resource exhausted" in low
        or "quota exceeded" in low
        or "rate limit" in low
        or "too many requests" in low
    )


def _looks_like_gemini_daily_quota(err_text: str) -> bool:
    low = (err_text or "").lower()
    return "requestsperday" in low or "perday" in low or "per day" in low


def _format_gemini_quota_error(err_text: str, *, model_name: str) -> str:
    model = (model_name or "").strip() or "gemini"
    limit = _gemini_quota_value(err_text)
    retry = _gemini_retry_after_seconds(err_text)
    is_daily = _looks_like_gemini_daily_quota(err_text)

    parts: List[str] = []
    if is_daily:
        head = f"Gemini 429: free-tier kunlik quota tugadi (model: {model})."
        if limit:
            head = f"Gemini 429: free-tier kunlik quota tugadi (limit: {limit}/kun, model: {model})."
        parts.append(head)
        parts.append("Yechim: ertaga qayta urinib ko'ring yoki billing/plan ni yoqing (quota oshadi).")
        parts.append("Tezda limitga tushmaslik: savol sonini kamaytiring (ayniqsa skan PDF).")
        return "\n".join(parts).strip()

    head = f"Gemini 429: so'rovlar limiti oshib ketdi (model: {model})."
    if limit:
        head = f"Gemini 429: so'rovlar limiti oshib ketdi (limit: {limit}, model: {model})."
    parts.append(head)
    if retry:
        parts.append(f"Yechim: {retry}s kutib qayta urinib ko'ring.")
    else:
        parts.append("Yechim: biroz kutib qayta urinib ko'ring yoki billing/plan ni yoqing.")
    return "\n".join(parts).strip()


def _looks_like_deadline_exceeded(err_text: str) -> bool:
    low = (err_text or "").lower()
    if not low:
        return False
    # Common Gemini/OpenAI/network timeout shapes.
    return (
        ("504" in low and ("deadline" in low or "timeout" in low))
        or "deadline exceeded" in low
        or "deadlineexceeded" in low
        or "deadlinexceeded" in low
        or "gateway timeout" in low
        or "upstream request timeout" in low
        or "context deadline exceeded" in low
        or "read timed out" in low
        or "connect timeout" in low
        or ("timed out" in low)
    )


def _format_deadline_error(
    err_text: str,
    *,
    provider: str,
    model_name: str = "",
    timeout_sec: int | None = None,
) -> str:
    prov = (provider or "ai").strip().lower()
    model = (model_name or "").strip()
    tmo = int(timeout_sec or 0)

    if prov == "gemini":
        head = f"Gemini 504: Deadline Exceeded (model: {model or 'gemini'})."
    elif prov == "openai":
        head = f"OpenAI timeout/504 (model: {model or 'openai'})."
    else:
        head = "AI 504: Deadline Exceeded."

    if tmo > 0:
        head += f" (timeout: {tmo}s)"

    parts: list[str] = [head]
    parts.append(
        "Yechim: savol sonini kamaytiring, mavzu/sahifa tanlang yoki matnni qisqartiring va qayta urinib ko'ring."
    )
    parts.append(
        "Tezroq ishlashi uchun: `.env` da `AI_MAX_TEXT_CHARS` ni kamaytiring (masalan 12000) "
        "yoki Gemini uchun `GEMINI_MODEL=gemini-flash-lite-latest` qilib ko'ring."
    )
    return "\n".join(parts).strip()


def _language_instruction(output_language: str) -> str:
    lang = (output_language or "source").strip().lower()
    if lang in {"source", "auto", "original"}:
        return (
            "Savollar va variantlar KIRISH matnining o'z tilida bo'lsin. Tarjima qilmang.\n"
            "Muhim: Agar kirish matni ruscha bo'lsa ruscha, inglizcha bo'lsa inglizcha yozing."
        )
    if lang in {"uz", "uzbek", "o'zbek", "ozbek"}:
        return "Savollar va variantlar O'zbek tilida bo'lsin."
    if lang in {"ru", "rus", "russian"}:
        return "Savollar va variantlar Rus tilida bo'lsin."
    if lang in {"en", "eng", "english"}:
        return "Savollar va variantlar Ingliz tilida bo'lsin."
    if lang in {"de", "deu", "german", "deutsch"}:
        return "Savollar va variantlar Nemis (German) tilida bo'lsin."
    if lang in {"tr", "tur", "turkish", "turkce"}:
        return "Savollar va variantlar Turk (Turkish) tilida bo'lsin."
    if lang in {"kk", "kazakh", "qazaq", "qazaqsha", "qozoq"}:
        return "Savollar va variantlar Qozoq (Kazakh) tilida bo'lsin."
    if lang in {"ar", "arab", "arabic"}:
        return "Savollar va variantlar Arab (Arabic) tilida bo'lsin."
    if lang in {"zh", "chi", "chinese", "zh-cn", "zh-hans", "mandarin"}:
        return "Savollar va variantlar Xitoy (Chinese) tilida bo'lsin."
    if lang in {"ko", "kor", "korean"}:
        return "Savollar va variantlar Koreys (Korean) tilida bo'lsin."
    raise AIServiceError("Til noto'g'ri. Tanlang: source | uz | ru | en | de | tr | kk | ar | zh | ko")



def _difficulty_instruction(difficulty: str) -> str:
    d = (difficulty or "mixed").strip().lower()
    if d in {"mixed", "mix", "aralash", "default", "random"}:
        return "Savollar turli qiyinlikda bo'lsin (oson, o'rta, qiyin)."
    if d in {"easy", "oson", "yengil", "beginner", "basic", "a1", "a2"}:
        return "Savollar asosan OSON bo'lsin (beginner)."
    if d in {"medium", "middle", "orta", "o'rta", "ortacha", "intermediate", "b1", "b2"}:
        return "Savollar asosan O'RTA darajada bo'lsin (intermediate)."
    if d in {"hard", "qiyin", "murakkab", "advanced", "c1", "c2"}:
        return "Savollar asosan QIYIN bo'lsin (advanced)."
    return "Savollar turli qiyinlikda bo'lsin (oson, o'rta, qiyin)."

def _topic_keywords(topic: str) -> List[str]:
    raw = (topic or "").strip().lower()
    if not raw:
        return []
    parts = [p for p in re.split(r"[^\w]+", raw) if p]
    # Drop very short tokens to avoid too many false matches (e.g., "va", "of", "the").
    parts = [p for p in parts if len(p) >= 3]
    # Preserve order but keep unique.
    seen = set()
    out: List[str] = []
    for p in parts:
        if p in seen:
            continue
        seen.add(p)
        out.append(p)
    return out[:12]


def _extract_relevant_text(text: str, topic: str, *, max_chars: int) -> str:
    """Best-effort extraction of topic-relevant snippets from large text."""
    cleaned = (text or "").strip()
    if not cleaned:
        return ""

    max_chars = max(1000, int(max_chars or 20000))

    keys = _topic_keywords(topic)
    if not keys:
        return cleaned[:max_chars]

    lower = cleaned.lower()
    spans: List[tuple[int, int]] = []

    # Collect a limited number of keyword hits to keep this fast.
    per_key_limit = 20
    for k in keys:
        if not k:
            continue
        for i, m in enumerate(re.finditer(re.escape(k), lower)):
            if i >= per_key_limit:
                break
            s = max(0, m.start() - 500)
            e = min(len(cleaned), m.end() + 800)
            spans.append((s, e))

    if not spans:
        return cleaned[:max_chars]

    spans.sort()
    merged: List[tuple[int, int]] = []
    for s, e in spans:
        if not merged:
            merged.append((s, e))
            continue
        ps, pe = merged[-1]
        if s <= pe + 200:
            merged[-1] = (ps, max(pe, e))
        else:
            merged.append((s, e))

    out_parts: List[str] = []
    total = 0
    for s, e in merged:
        piece = cleaned[s:e].strip()
        if not piece:
            continue
        sep = "\n...\n" if out_parts else ""
        add_len = len(sep) + len(piece)
        if total + add_len > max_chars:
            remaining = max_chars - total - len(sep)
            if remaining > 200:
                out_parts.append(sep + piece[:remaining])
            break
        out_parts.append(sep + piece)
        total += add_len

    result = "".join(out_parts).strip()
    if len(result) < 500:
        return cleaned[:max_chars]
    return result[:max_chars]


def extract_text_from_file(file_path: str) -> str:
    suffix = Path(file_path).suffix.lower()
    char_limit = int(os.getenv("EXTRACT_CHAR_LIMIT", "200000"))

    if suffix == ".pdf":
        return _extract_text_from_pdf(file_path, char_limit=char_limit)
    if suffix == ".docx":
        return _extract_text_from_docx(file_path, char_limit=char_limit)
    if suffix == ".pptx":
        return _extract_text_from_pptx(file_path, char_limit=char_limit)
    if suffix in {".txt", ".md"}:
        return Path(file_path).read_text(encoding="utf-8", errors="ignore")[:char_limit]

    raise AIServiceError("Unsupported file type. Send .pdf, .docx, .pptx, .txt, or .md")


def _extract_text_from_pdf(file_path: str, *, char_limit: int) -> str:
    try:
        import fitz  # PyMuPDF
    except ImportError as exc:
        raise AIServiceError("PyMuPDF not installed. Install: pip install PyMuPDF") from exc

    parts: List[str] = []
    total = 0
    with fitz.open(file_path) as doc:
        for page in doc:
            page_text = page.get_text()
            if not page_text:
                continue
            parts.append(page_text)
            total += len(page_text)
            if total >= char_limit:
                break
    return "\n".join(parts)


def _extract_text_from_docx(file_path: str, *, char_limit: int) -> str:
    try:
        import docx  # python-docx
    except ImportError as exc:
        raise AIServiceError("python-docx not installed. Install: pip install python-docx") from exc

    doc = docx.Document(file_path)
    parts: List[str] = []
    total = 0
    for p in doc.paragraphs:
        t = p.text or ""
        if not t:
            continue
        parts.append(t)
        total += len(t)
        if total >= char_limit:
            break
    return "\n".join(parts)


def _extract_text_from_pptx(file_path: str, *, char_limit: int) -> str:
    """Extract text from a .pptx.

    Prefer python-pptx when available; fall back to stdlib zip+xml extraction so the
    bot still works even if python-pptx isn't installed.
    """

    try:
        from pptx import Presentation  # python-pptx
    except ImportError:
        import zipfile
        import xml.etree.ElementTree as ET

        def _slide_num(name: str) -> int:
            m = re.search(r"slide(\d+)\.xml$", name)
            return int(m.group(1)) if m else 0

        parts: List[str] = []
        total = 0
        with zipfile.ZipFile(file_path) as z:
            slides = [
                n
                for n in z.namelist()
                if n.startswith("ppt/slides/slide") and n.endswith(".xml")
            ]
            slides.sort(key=_slide_num)

            for name in slides:
                try:
                    root = ET.fromstring(z.read(name))
                except Exception:
                    continue

                texts: List[str] = []
                for el in root.iter():
                    # <a:t> nodes hold text runs.
                    if el.tag.endswith("}t") and el.text:
                        t = el.text.strip()
                        if t:
                            texts.append(t)

                if not texts:
                    continue

                chunk = " ".join(texts)
                parts.append(chunk)
                total += len(chunk)
                if total >= char_limit:
                    break

        return "\n".join(parts)[:char_limit]

    prs = Presentation(file_path)
    parts: List[str] = []
    total = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            # Most text-bearing shapes expose .text
            t = ""
            try:
                t = str(getattr(shape, "text", "") or "")
            except Exception:
                t = ""
            if t:
                parts.append(t)
                total += len(t)
                if total >= char_limit:
                    return "\n".join(parts)

            # Tables need explicit cell extraction
            try:
                has_table = bool(getattr(shape, "has_table", False))
            except Exception:
                has_table = False
            if not has_table:
                continue
            try:
                table = shape.table
            except Exception:
                continue
            for row in table.rows:
                for cell in row.cells:
                    ct = str(getattr(cell, "text", "") or "")
                    if not ct:
                        continue
                    parts.append(ct)
                    total += len(ct)
                    if total >= char_limit:
                        return "\n".join(parts)

    return "\n".join(parts)[:char_limit]


def _load_json_from_text(text: str) -> Any:
    raw = (text or "").strip()
    if not raw:
        raise AIServiceError("Empty AI response")

    # Strip common markdown fences.
    raw = re.sub(r"^```(?:json)?\s*", "", raw, flags=re.IGNORECASE)
    raw = re.sub(r"\s*```$", "", raw)

    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        # Fallback: attempt to parse the first JSON object/array inside the text.
        match = re.search(r"(\{.*\}|\[.*\])", raw, flags=re.S)
        if not match:
            raise AIServiceError("AI response is not valid JSON")
        try:
            return json.loads(match.group(1))
        except json.JSONDecodeError as exc:
            raise AIServiceError("AI response JSON parsing failed") from exc


def _normalize_quiz(data: Any) -> List[Dict[str, Any]]:
    if isinstance(data, dict):
        for key in ("quiz", "questions", "items"):
            if key in data and isinstance(data[key], list):
                data = data[key]
                break

    if not isinstance(data, list):
        raise AIServiceError("AI JSON must be a list or an object with a 'quiz' list")

    out: List[Dict[str, Any]] = []
    for item in data:
        if not isinstance(item, dict):
            continue

        question = str(item.get("question") or item.get("text") or "").strip()
        options = item.get("options") or item.get("variants") or item.get("answers") or []
        explanation = str(item.get("explanation") or item.get("comment") or "").strip()

        if isinstance(options, str):
            options = [s.strip() for s in re.split(r"[|\n]", options) if s.strip()]

        if not isinstance(options, list):
            continue

        options = [str(o).strip() for o in options if str(o).strip()]
        if len(options) < 4:
            continue
        if len(options) > 4:
            options = options[:4]

        correct = item.get("correct_index")
        if correct is None:
            correct = item.get("correct") or item.get("answer_index")

        correct_index = 0
        try:
            correct_index = int(correct)
        except Exception:
            if isinstance(correct, str):
                letter = correct.strip().upper()
                mapping = {"A": 0, "B": 1, "C": 2, "D": 3}
                if letter in mapping:
                    correct_index = mapping[letter]

        if correct_index < 0 or correct_index > 3:
            correct_index = max(0, min(3, correct_index))

        if not question:
            continue

        # Shuffle options so the correct answer isn't always A (0).
        try:
            perm = [0, 1, 2, 3]
            random.shuffle(perm)
            options = [options[i] for i in perm]
            correct_index = int(perm.index(int(correct_index)))
        except Exception:
            pass

        out.append(
            {
                "question": question,
                "options": options,
                "correct_index": correct_index,
                "explanation": explanation,
            }
        )

    if not out:
        raise AIServiceError("AI did not return any valid questions")
    return out


def _normalize_receipt_review(data: Any) -> Dict[str, Any]:
    """Normalize AI receipt review JSON into a stable shape.

    Expected output:
      {"verdict":"approve|suspicious", "confidence":0..1, "amount_uzs":int, "reason":str}
    """

    if not isinstance(data, dict):
        raise AIServiceError("Receipt review JSON must be an object")

    verdict = str(data.get("verdict") or data.get("status") or data.get("decision") or "").strip().lower()
    if verdict not in {"approve", "suspicious", "reject"}:
        verdict = "suspicious"
    if verdict == "reject":
        # We don't auto-reject payments; treat as suspicious.
        verdict = "suspicious"

    conf_raw = data.get("confidence")
    if conf_raw is None:
        conf_raw = data.get("conf") or data.get("score")
    try:
        confidence = float(conf_raw)
    except Exception:
        confidence = 0.0
    confidence = max(0.0, min(1.0, confidence))

    amount_raw = data.get("amount_uzs")
    if amount_raw is None:
        amount_raw = data.get("amount") or data.get("sum") or data.get("uzs")
    amount_uzs = 0
    try:
        if isinstance(amount_raw, str):
            cleaned = re.sub(r"[^\\d]", "", amount_raw)
            amount_uzs = int(cleaned) if cleaned else 0
        else:
            amount_uzs = int(amount_raw or 0)
    except Exception:
        amount_uzs = 0
    amount_uzs = max(0, int(amount_uzs))

    reason = str(data.get("reason") or data.get("comment") or data.get("explanation") or "").strip()
    if not reason:
        reason = "no_reason"

    return {
        "verdict": verdict,
        "confidence": confidence,
        "amount_uzs": amount_uzs,
        "reason": reason,
    }


@dataclass
class AIService:
    provider: str = AI_PROVIDER
    openai_api_key: Optional[str] = OPENAI_API_KEY
    gemini_api_key: Optional[str] = GEMINI_API_KEY
    openai_model: str = OPENAI_MODEL
    gemini_model: str = GEMINI_MODEL

    def _pick_provider(self) -> str:
        provider = (self.provider or "auto").strip().lower()
        if provider not in {"auto", "openai", "gemini"}:
            raise AIServiceError("AI_PROVIDER faqat: auto | openai | gemini")

        if provider == "openai":
            if not (self.openai_api_key or "").strip():
                raise AIServiceError("OPENAI_API_KEY yo'q. .env ga qo'shing yoki AI_PROVIDER=gemini qiling.")
            if _looks_like_gemini_key(self.openai_api_key):
                raise AIServiceError(
                    "OPENAI_API_KEY ga Gemini kaliti (AIza...) yozilib qolgan.\n"
                    "Yechim: .env da GEMINI_API_KEY=AIza... va AI_PROVIDER=gemini qiling "
                    "(yoki OPENAI_API_KEY ni haqiqiy OpenAI kalit (sk-...) ga almashtiring)."
                )
            if not _looks_like_openai_key(self.openai_api_key):
                raise AIServiceError("OPENAI_API_KEY formati noto'g'ri (odatda sk-... bilan boshlanadi).")
            return "openai"

        if provider == "gemini":
            if not (self.gemini_api_key or "").strip():
                raise AIServiceError("GEMINI_API_KEY yo'q. .env ga qo'shing yoki AI_PROVIDER=openai qiling.")
            if _looks_like_openai_key(self.gemini_api_key):
                raise AIServiceError("GEMINI_API_KEY ga OpenAI kaliti (sk-...) yozilib qolgan.")
            return "gemini"

        # auto: choose by presence + key shape to avoid common misconfigurations
        if (self.openai_api_key or "").strip() and _looks_like_openai_key(self.openai_api_key):
            return "openai"
        if (self.gemini_api_key or "").strip():
            return "gemini"
        if _looks_like_gemini_key(self.openai_api_key) and not (self.gemini_api_key or "").strip():
            raise AIServiceError(
                "Sizda faqat OPENAI_API_KEY bor, lekin u Gemini kalitiga (AIza...) o'xshaydi.\n"
                "Yechim: OPENAI_API_KEY ni olib tashlang, GEMINI_API_KEY ga qo'ying va AI_PROVIDER=gemini qiling."
            )
        if (self.openai_api_key or "").strip():
            raise AIServiceError("OPENAI_API_KEY bor, lekin OpenAI kalitiga o'xshamaydi (sk-...).")
        raise AIServiceError("OPENAI_API_KEY yoki GEMINI_API_KEY ni .env da sozlang.")

    async def generate_quiz_from_text(
        self,
        text: str,
        question_count: int = 5,
        output_language: str = "source",
        difficulty: str = "mixed",
        focus_topic: str = "",
    ) -> List[Dict[str, Any]]:
        provider = self._pick_provider()

        cleaned = (text or "").strip()
        if len(cleaned) < 200:
            raise AIServiceError("Text is too short to generate a quiz")

        lang_instruction = _language_instruction(output_language)

        diff_extra = _difficulty_instruction(difficulty).strip()

        # Keep prompts bounded for speed/cost. Override via env if needed.
        max_chars = int(os.getenv("AI_MAX_TEXT_CHARS", "12000"))
        max_chars = max(5000, min(120000, max_chars))
        focus_topic = (focus_topic or "").strip()
        if focus_topic:
            cleaned = _extract_relevant_text(cleaned, focus_topic, max_chars=max_chars)
        else:
            cleaned = cleaned[:max_chars]

        topic_extra = ""
        if focus_topic:
            topic_extra = (
                "Faqat quyidagi mavzuga oid savollar tuzing va matndan faqat shu mavzuga tegishli "
                f"qismlarni ishlating: {focus_topic}"
            )

        # Difficulty
        extra0 = ""
        if topic_extra:
            extra0 = topic_extra
        if diff_extra:
            extra0 = (diff_extra + ("\n\n" + extra0 if extra0 else "")).strip()

        batch_size = int(os.getenv("AI_BATCH_SIZE", "10") or 10)
        batch_size = max(1, min(20, batch_size))

        if provider == "openai":
            questions = await self._generate_openai(
                cleaned,
                min(question_count, batch_size),
                lang_instruction=lang_instruction,
                extra_instructions=extra0,
            )
        elif provider == "gemini":
            questions = await self._generate_gemini(
                cleaned,
                min(question_count, batch_size),
                lang_instruction=lang_instruction,
                extra_instructions=extra0,
            )
        else:
            raise AIServiceError(f"Unknown AI provider: {provider}")

        uniq: List[Dict[str, Any]] = []
        seen = set()
        for q in questions:
            key = str(q.get("question") or "").strip().lower()
            if not key or key in seen:
                continue
            seen.add(key)
            uniq.append(q)

        # If the model returned fewer than requested, try to fill the gap (bounded retries).
        max_retries = int(os.getenv("AI_FILL_RETRIES", "2"))
        max_retries = max(0, min(5, max_retries))
        max_iters = max_retries + (question_count + batch_size - 1) // batch_size
        tries = 0
        while len(uniq) < question_count and tries < max_iters:
            tries += 1
            remaining = question_count - len(uniq)
            ask = min(remaining, batch_size)
            avoid = "\n".join(f"- {q['question']}" for q in uniq[:12] if q.get("question"))
            extra = "Quyidagi savollarni takrorlamang:\n" + avoid if avoid else ""

            if topic_extra:
                extra = (topic_extra + ("\n\n" + extra if extra else "")).strip()
            if diff_extra:
                extra = (diff_extra + ("\n\n" + extra if extra else "")).strip()

            if provider == "openai":
                more = await self._generate_openai(
                    cleaned,
                    ask,
                    lang_instruction=lang_instruction,
                    extra_instructions=extra,
                )
            else:
                more = await self._generate_gemini(
                    cleaned,
                    ask,
                    lang_instruction=lang_instruction,
                    extra_instructions=extra,
                )

            for q in more:
                key = str(q.get("question") or "").strip().lower()
                if not key or key in seen:
                    continue
                seen.add(key)
                uniq.append(q)

        return uniq[:question_count]

    async def generate_quiz_from_topic(
        self,
        topic: str,
        question_count: int = 5,
        output_language: str = "source",
        difficulty: str = "mixed",
    ) -> List[Dict[str, Any]]:
        provider = self._pick_provider()
        topic = (topic or "").strip()
        if len(topic) < 3:
            raise AIServiceError("Mavzu juda qisqa. Iltimos, aniqroq yozing.")

        lang_instruction = _language_instruction(output_language)
        diff_extra_topic = _difficulty_instruction(difficulty).strip()

        batch_size = int(os.getenv("AI_BATCH_SIZE", "10") or 10)
        batch_size = max(1, min(20, batch_size))

        if provider == "openai":
            questions = await self._generate_openai_topic(
                topic,
                min(question_count, batch_size),
                lang_instruction=lang_instruction,
                extra_instructions=diff_extra_topic,
            )
        elif provider == "gemini":
            questions = await self._generate_gemini_topic(
                topic,
                min(question_count, batch_size),
                lang_instruction=lang_instruction,
                extra_instructions=diff_extra_topic,
            )
        else:
            raise AIServiceError(f"Unknown AI provider: {provider}")

        uniq: List[Dict[str, Any]] = []
        seen = set()
        for q in questions:
            key = str(q.get("question") or "").strip().lower()
            if not key or key in seen:
                continue
            seen.add(key)
            uniq.append(q)

        max_retries = int(os.getenv("AI_FILL_RETRIES", "4") or 4)
        max_retries = max(0, min(8, max_retries))
        max_iters = max(
            max_retries + (question_count + batch_size - 1) // batch_size,
            (question_count * 2 + batch_size - 1) // batch_size,
        )
        tries = 0
        while len(uniq) < question_count and tries < max_iters:
            tries += 1
            remaining = question_count - len(uniq)
            if remaining <= 0:
                break

            parallel = int(os.getenv("AI_PARALLEL_BATCHES", "2") or 2)
            parallel = max(1, min(5, parallel))

            avoid = "\n".join(f"- {q['question']}" for q in uniq[:12] if q.get("question"))
            extra = "Quyidagi savollarni takrorlamang:\n" + avoid if avoid else ""
            if diff_extra_topic:
                extra = (diff_extra_topic + ("\n\n" + extra if extra else "")).strip()

            tasks = []
            to_ask = remaining
            for _ in range(parallel):
                if to_ask <= 0:
                    break
                ask = min(to_ask, batch_size)
                to_ask -= ask
                if provider == "openai":
                    tasks.append(
                        self._generate_openai_topic(
                            topic,
                            ask,
                            lang_instruction=lang_instruction,
                            extra_instructions=extra,
                        )
                    )
                else:
                    tasks.append(
                        self._generate_gemini_topic(
                            topic,
                            ask,
                            lang_instruction=lang_instruction,
                            extra_instructions=extra,
                        )
                    )

            results = await asyncio.gather(*tasks, return_exceptions=True)
            for res in results:
                if isinstance(res, Exception):
                    raise res
                for q in res:
                    key = str(q.get("question") or "").strip().lower()
                    if not key or key in seen:
                        continue
                    seen.add(key)
                    uniq.append(q)
        return uniq[:question_count]

    async def generate_quiz_from_images(
        self,
        image_paths: List[str],
        *,
        output_language: str = "source",
    ) -> List[Dict[str, Any]]:
        """Generate exactly 1 question per image (Gemini vision)."""
        provider = self._pick_provider()
        if provider != "gemini":
            raise AIServiceError("Rasmli savollar hozircha faqat Gemini bilan ishlaydi (AI_PROVIDER=gemini).")

        paths = [str(p) for p in (image_paths or []) if str(p).strip()]
        if not paths:
            raise AIServiceError("Image list is empty")

        lang_instruction = _language_instruction(output_language)

        out: List[Dict[str, Any]] = []
        for p in paths:
            q = await self._generate_gemini_image(p, lang_instruction=lang_instruction)
            out.append(q)
        return out


    async def review_payment_receipt_text(self, receipt_text: str, *, expected_amount_uzs: int) -> Dict[str, Any]:
        """AI check for payment receipt shared as plain text (share-check).

        Returns normalized dict:
          {"verdict":"approve|suspicious", "confidence":0..1, "amount_uzs":int, "reason":str}
        """

        expected_amount_uzs = int(expected_amount_uzs or 0)
        text = (receipt_text or "").strip()
        if len(text) < 10:
            return {"verdict": "suspicious", "confidence": 0.0, "amount_uzs": 0, "reason": "text_too_short"}

        provider = self._pick_provider()

        timeout_sec = float(os.getenv("RECEIPT_AI_TIMEOUT_SEC", os.getenv("AI_REQUEST_TIMEOUT_SEC", "40")) or 40)
        timeout_sec = max(5.0, min(120.0, timeout_sec))

        system_prompt = (
            "You are a payment receipt verifier.\\n"
            "Given a receipt (TEXT) you must decide if it looks like a completed payment for EXACT amount.\\n"
            f"Expected amount: {expected_amount_uzs} UZS.\\n"
            "Rules:\\n"
            "- If amount is missing or not equal to expected, verdict MUST be suspicious.\\n"
            "- If the receipt indicates pending/failed/canceled, verdict MUST be suspicious.\\n"
            "- If the text is too short / unclear, verdict MUST be suspicious.\\n"
            "Output ONLY JSON object:\\n"
            '{"verdict":"approve|suspicious", "confidence":0.0, "amount_uzs":0, "reason":"..."}'
        )

        if provider == "openai":
            try:
                from openai import AsyncOpenAI
            except ImportError as exc:
                raise AIServiceError("openai not installed. Install: pip install openai") from exc
            if not self.openai_api_key:
                raise AIServiceError("OPENAI_API_KEY is not set")

            client = AsyncOpenAI(api_key=self.openai_api_key)
            total_timeout = float(os.getenv("OPENAI_TOTAL_TIMEOUT_SEC", "0") or 0)
            if total_timeout <= 0:
                total_timeout = min(120.0, timeout_sec + 20.0)

            try:
                resp = await asyncio.wait_for(
                    client.chat.completions.create(
                        model=self.openai_model,
                        messages=[
                            {"role": "system", "content": system_prompt},
                            {"role": "user", "content": text},
                        ],
                        temperature=0.0,
                        response_format={"type": "json_object"},
                    ),
                    timeout=total_timeout,
                )
            except asyncio.TimeoutError as exc:
                raise AIServiceError(
                    _format_deadline_error("timeout", provider="openai", model_name=str(self.openai_model), timeout_sec=int(total_timeout))
                ) from exc
            except Exception as exc:
                msg = str(exc)
                if "invalid_api_key" in msg or "Incorrect API key" in msg or "Error code: 401" in msg:
                    raise AIServiceError("OpenAI 401: API key noto'g'ri yoki bekor qilingan.") from exc
                if _looks_like_deadline_exceeded(msg):
                    raise AIServiceError(
                        _format_deadline_error(msg, provider="openai", model_name=str(self.openai_model), timeout_sec=int(total_timeout))
                    ) from exc
                raise

            content = resp.choices[0].message.content or ""
            data = _load_json_from_text(content)
            return _normalize_receipt_review(data)

        # gemini
        try:
            import google.generativeai as genai
        except ImportError as exc:
            raise AIServiceError("google-generativeai not installed. Install: pip install google-generativeai") from exc

        if not self.gemini_api_key:
            raise AIServiceError("GEMINI_API_KEY is not set")

        genai.configure(api_key=self.gemini_api_key)
        model_name = (self.gemini_model or "").strip()
        if not model_name:
            raise AIServiceError("GEMINI_MODEL bo'sh. .env da GEMINI_MODEL=... qo'ying.")
        if model_name.startswith("models/"):
            model_name = model_name.split("/", 1)[1]

        model = genai.GenerativeModel(model_name)
        prompt = f"{system_prompt}\\n\\nRECEIPT_TEXT:\\n{text}".strip()

        def _call() -> str:
            try:
                try:
                    r = model.generate_content(
                        prompt,
                        generation_config={
                            "temperature": 0.0,
                            "response_mime_type": "application/json",
                        },
                    )
                except TypeError:
                    r = model.generate_content(prompt)
                return getattr(r, "text", "") or ""
            except Exception as exc:
                msg = str(exc)
                if _looks_like_gemini_quota_error(msg):
                    raise AIServiceError(_format_gemini_quota_error(msg, model_name=model_name)) from exc
                if _looks_like_gemini_auth_error(msg):
                    raise AIServiceError(_format_gemini_auth_error(msg)) from exc
                if _looks_like_deadline_exceeded(msg):
                    raise AIServiceError(
                        _format_deadline_error(msg, provider="gemini", model_name=str(model_name), timeout_sec=int(timeout_sec))
                    ) from exc
                raise

        try:
            raw = await asyncio.wait_for(asyncio.to_thread(_call), timeout=timeout_sec)
        except asyncio.TimeoutError as exc:
            raise AIServiceError(
                _format_deadline_error("timeout", provider="gemini", model_name=str(model_name), timeout_sec=int(timeout_sec))
            ) from exc

        data = _load_json_from_text(raw)
        return _normalize_receipt_review(data)

    async def review_payment_receipt_image(self, image_path: str, *, expected_amount_uzs: int) -> Dict[str, Any]:
        """AI check for payment receipt screenshot (image/PDF rendered to image).

        Uses Gemini vision (requires GEMINI_API_KEY).
        """

        expected_amount_uzs = int(expected_amount_uzs or 0)
        p = str(image_path or "").strip()
        if not p:
            raise AIServiceError("image_path is empty")
        if not Path(p).exists():
            raise AIServiceError("Receipt image not found")

        # Even if AI_PROVIDER=openai, receipt screenshots need Gemini vision.
        if not (self.gemini_api_key or "").strip():
            raise AIServiceError("GEMINI_API_KEY is required for receipt screenshot verification")

        timeout_sec = float(os.getenv("RECEIPT_AI_TIMEOUT_SEC", os.getenv("AI_REQUEST_TIMEOUT_SEC", "40")) or 40)
        timeout_sec = max(5.0, min(120.0, timeout_sec))

        prompt = (
            "You are a payment receipt verifier.\\n"
            "Given a receipt SCREENSHOT/IMAGE you must decide if it looks like a completed payment for EXACT amount.\\n"
            f"Expected amount: {expected_amount_uzs} UZS.\\n"
            "Rules:\\n"
            "- If amount is missing or not equal to expected, verdict MUST be suspicious.\\n"
            "- If the receipt indicates pending/failed/canceled, verdict MUST be suspicious.\\n"
            "- If the image is not a payment receipt or is unclear, verdict MUST be suspicious.\\n"
            "Output ONLY JSON object:\\n"
            '{"verdict":"approve|suspicious", "confidence":0.0, "amount_uzs":0, "reason":"..."}'
        )

        try:
            import google.generativeai as genai
        except ImportError as exc:
            raise AIServiceError("google-generativeai not installed. Install: pip install google-generativeai") from exc

        genai.configure(api_key=self.gemini_api_key)
        model_name = (self.gemini_model or "").strip()
        if not model_name:
            raise AIServiceError("GEMINI_MODEL bo'sh. .env da GEMINI_MODEL=... qo'ying.")
        if model_name.startswith("models/"):
            model_name = model_name.split("/", 1)[1]

        model = genai.GenerativeModel(model_name)

        def _call() -> str:
            uploaded = None
            try:
                uploaded = genai.upload_file(p)
                try:
                    r = model.generate_content(
                        [prompt, uploaded],
                        generation_config={
                            "temperature": 0.0,
                            "response_mime_type": "application/json",
                        },
                    )
                except TypeError:
                    r = model.generate_content([prompt, uploaded])
                return getattr(r, "text", "") or ""
            finally:
                if uploaded is not None:
                    try:
                        genai.delete_file(uploaded)
                    except Exception:
                        pass

        try:
            raw = await asyncio.wait_for(asyncio.to_thread(_call), timeout=timeout_sec)
        except asyncio.TimeoutError as exc:
            raise AIServiceError(
                _format_deadline_error("timeout", provider="gemini", model_name=str(model_name), timeout_sec=int(timeout_sec))
            ) from exc
        except Exception as exc:
            msg = str(exc)
            if _looks_like_gemini_quota_error(msg):
                raise AIServiceError(_format_gemini_quota_error(msg, model_name=model_name)) from exc
            if _looks_like_gemini_auth_error(msg):
                raise AIServiceError(_format_gemini_auth_error(msg)) from exc
            if _looks_like_deadline_exceeded(msg):
                raise AIServiceError(
                    _format_deadline_error(msg, provider="gemini", model_name=str(model_name), timeout_sec=int(timeout_sec))
                ) from exc
            raise

        data = _load_json_from_text(raw)
        return _normalize_receipt_review(data)

    async def _generate_openai(
        self,
        text: str,
        question_count: int,
        *,
        lang_instruction: str,
        extra_instructions: str = "",
    ) -> List[Dict[str, Any]]:
        try:
            from openai import AsyncOpenAI
        except ImportError as exc:
            raise AIServiceError("openai not installed. Install: pip install openai") from exc

        if not self.openai_api_key:
            raise AIServiceError("OPENAI_API_KEY is not set")

        client = AsyncOpenAI(api_key=self.openai_api_key)

        timeout_sec = float(os.getenv("OPENAI_TIMEOUT_SEC", os.getenv("AI_REQUEST_TIMEOUT_SEC", "60")) or 60)
        timeout_sec = max(5.0, min(600.0, timeout_sec))
        total_timeout = float(os.getenv("OPENAI_TOTAL_TIMEOUT_SEC", "0") or 0)
        if total_timeout <= 0:
            total_timeout = min(900.0, timeout_sec + 30.0)

        system_prompt = (
            "Sen professional testologsan. Berilgan matn asosida ko'p tanlovli test yarat.\n"
            f"- Savollar soni: {question_count}\n"
            "- Har savolda 4 ta variant bo'lsin\n"
            "- Variantlar bir xil uslubda bo'lsin (hammasi ibora yoki hammasi 1 gap)\n"
            "- Variantlar uzunligi bir-biriga yaqin bo'lsin: eng uzun va eng qisqa variant farqi 2-3 so'zdan oshmasin\n"
            "- Juda qisqa (1-2 so'zli) va juda uzun variantlardan qoching; har bir variant taxminan 4-10 so'z bo'lsin\n"
            "- Variantlarda izoh/tushuntirish bo'lmasin; izoh faqat explanation ga yozilsin\n"
            "- Noto'g'ri variantlar ham mantiqan ishonarli bo'lsin, to'g'ri javob ko'zga tashlanib qolmasin\n"
            "- correct_index 0..3 bo'lsin\n"
            "- explanation qisqa bo'lsin\n"
            f"- Til: {lang_instruction}\n"
            "Natijani faqat JSON ko'rinishida qaytar: {\"quiz\": [...]}.\n"
            "Har element: {\"question\": \"...\", \"options\": [\"A\",\"B\",\"C\",\"D\"], "
            "\"correct_index\": 0, \"explanation\": \"...\"}"
        )

        user_prompt = f"{extra_instructions}\n\nMatn:\n{text}".strip()

        try:
            response = await asyncio.wait_for(
                client.chat.completions.create(
                    model=self.openai_model,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_prompt},
                    ],
                    temperature=0.2,
                    response_format={"type": "json_object"},
                ),
                timeout=total_timeout,
            )
        except asyncio.TimeoutError as exc:
            raise AIServiceError(
                _format_deadline_error("timeout", provider="openai", model_name=str(self.openai_model), timeout_sec=int(total_timeout))
            ) from exc
        except Exception as exc:
            msg = str(exc)
            if "invalid_api_key" in msg or "Incorrect API key" in msg or "Error code: 401" in msg:
                raise AIServiceError(
                    "OpenAI 401: API key noto'g'ri yoki bekor qilingan. .env da OPENAI_API_KEY ni tekshiring."
                ) from exc
            if _looks_like_deadline_exceeded(msg):
                raise AIServiceError(
                    _format_deadline_error(msg, provider="openai", model_name=str(self.openai_model), timeout_sec=int(total_timeout))
                ) from exc
            raise

        content = response.choices[0].message.content or ""
        data = _load_json_from_text(content)
        return _normalize_quiz(data)

    async def _generate_openai_topic(
        self,
        topic: str,
        question_count: int,
        *,
        lang_instruction: str,
        extra_instructions: str = "",
    ) -> List[Dict[str, Any]]:
        try:
            from openai import AsyncOpenAI
        except ImportError as exc:
            raise AIServiceError("openai not installed. Install: pip install openai") from exc

        if not self.openai_api_key:
            raise AIServiceError("OPENAI_API_KEY is not set")

        client = AsyncOpenAI(api_key=self.openai_api_key)

        timeout_sec = float(os.getenv("OPENAI_TIMEOUT_SEC", os.getenv("AI_REQUEST_TIMEOUT_SEC", "60")) or 60)
        timeout_sec = max(5.0, min(600.0, timeout_sec))
        total_timeout = float(os.getenv("OPENAI_TOTAL_TIMEOUT_SEC", "0") or 0)
        if total_timeout <= 0:
            total_timeout = min(900.0, timeout_sec + 30.0)
        system_prompt = f"""Sen professional testologsan. Berilgan MAVZU bo'yicha ko'p tanlovli test yarat.
- Savollar soni: {question_count}
- Mavzuni KENG qamrab ol: asosiy tushunchalar, sabab-oqibat, taqqoslash, amaliy qo'llanish, tipik xatolar, terminlar/ta'riflar, klassifikatsiya (mavzuga mos bo'lsa).
- Savollar bir xil qolipda bo'lmasin: savolning boshlanishi va uslubi turlicha bo'lsin (hammasi bir xil so'z bilan boshlanmasin).
- Mavzu so'zi bilan har doim boshlamang. (Masalan, 'Olma ...' deb ketma-ket boshlamang.)
- Har savolda 4 ta variant bo'lsin
- Variantlar bir xil uslubda bo'lsin (hammasi ibora yoki hammasi 1 gap)
- Variantlar uzunligi bir-biriga yaqin bo'lsin: eng uzun va eng qisqa variant farqi 2-3 so'zdan oshmasin
- Juda qisqa (1-2 so'zli) va juda uzun variantlardan qoching; har bir variant taxminan 4-10 so'z bo'lsin
- Variantlarda izoh/tushuntirish bo'lmasin; izoh faqat explanation ga yozilsin
- Noto'g'ri variantlar ham mantiqan ishonarli bo'lsin, to'g'ri javob ko'zga tashlanib qolmasin
- correct_index 0..3 bo'lsin
- explanation qisqa bo'lsin
- Til: {lang_instruction}
- Muhim: EXACTLY shu miqdorda savol qaytar (ro'yxat uzunligi aynan Savollar soni bo'lsin).
Natijani faqat JSON ko'rinishida qaytar: {{\"quiz\": [...]}}.
Har element: {{\"question\": \"...\", \"options\": [\"A\",\"B\",\"C\",\"D\"], \"correct_index\": 0, \"explanation\": \"...\"}}"""

        user_prompt = f"{extra_instructions}\n\nMavzu: {topic}".strip()

        try:
            response = await asyncio.wait_for(
                client.chat.completions.create(
                    model=self.openai_model,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_prompt},
                    ],
                    temperature=topic_temp,
                    response_format={"type": "json_object"},
                ),
                timeout=total_timeout,
            )
        except asyncio.TimeoutError as exc:
            raise AIServiceError(
                _format_deadline_error("timeout", provider="openai", model_name=str(self.openai_model), timeout_sec=int(total_timeout))
            ) from exc
        except Exception as exc:
            msg = str(exc)
            if "invalid_api_key" in msg or "Incorrect API key" in msg or "Error code: 401" in msg:
                raise AIServiceError(
                    "OpenAI 401: API key noto'g'ri yoki bekor qilingan. .env da OPENAI_API_KEY ni tekshiring."
                ) from exc
            if _looks_like_deadline_exceeded(msg):
                raise AIServiceError(
                    _format_deadline_error(msg, provider="openai", model_name=str(self.openai_model), timeout_sec=int(total_timeout))
                ) from exc
            raise

        content = response.choices[0].message.content or ""
        data = _load_json_from_text(content)
        return _normalize_quiz(data)

    async def _generate_gemini(
        self,
        text: str,
        question_count: int,
        *,
        lang_instruction: str,
        extra_instructions: str = "",
    ) -> List[Dict[str, Any]]:
        try:
            import google.generativeai as genai
        except ImportError as exc:
            raise AIServiceError(
                "google-generativeai not installed. Install: pip install google-generativeai"
            ) from exc

        if not self.gemini_api_key:
            raise AIServiceError("GEMINI_API_KEY is not set")

        genai.configure(api_key=self.gemini_api_key)
        model_name = (self.gemini_model or "").strip()
        if not model_name:
            raise AIServiceError("GEMINI_MODEL bo'sh. .env da GEMINI_MODEL=... qo'ying.")

        # list_models() usually returns "models/...", but GenerativeModel typically accepts
        # plain names too; normalize so users can paste either form into .env.
        if model_name.startswith("models/"):
            model_name = model_name.split("/", 1)[1]

        model = genai.GenerativeModel(model_name)

        timeout_sec = int(os.getenv("GEMINI_TIMEOUT_SEC", os.getenv("AI_REQUEST_TIMEOUT_SEC", "60")) or 60)
        timeout_sec = max(5, min(300, timeout_sec))
        total_timeout = int(os.getenv("GEMINI_TOTAL_TIMEOUT_SEC", str(timeout_sec + 30)) or (timeout_sec + 30))
        total_timeout = max(timeout_sec, min(600, total_timeout))

        prompt = (
            "Sen professional testologsan. Berilgan matn asosida ko'p tanlovli test yarat.\n"
            f"- Savollar soni: {question_count}\n"
            "- Har savolda 4 ta variant bo'lsin\n"
            "- Variantlar bir xil uslubda bo'lsin (hammasi ibora yoki hammasi 1 gap)\n"
            "- Variantlar uzunligi bir-biriga yaqin bo'lsin: eng uzun va eng qisqa variant farqi 2-3 so'zdan oshmasin\n"
            "- Juda qisqa (1-2 so'zli) va juda uzun variantlardan qoching; har bir variant taxminan 4-10 so'z bo'lsin\n"
            "- Variantlarda izoh/tushuntirish bo'lmasin; izoh faqat explanation ga yozilsin\n"
            "- Noto'g'ri variantlar ham mantiqan ishonarli bo'lsin, to'g'ri javob ko'zga tashlanib qolmasin\n"
            "- correct_index 0..3 bo'lsin\n"
            "- explanation qisqa bo'lsin\n"
            f"- Til: {lang_instruction}\n"
            "Natijani faqat JSON ko'rinishida qaytar: {\"quiz\": [...]}.\n"
            "Har element: {\"question\": \"...\", \"options\": [\"A\",\"B\",\"C\",\"D\"], "
            "\"correct_index\": 0, \"explanation\": \"...\"}\n\n"
            f"{extra_instructions}\n\nMatn:\n{text}".strip()
        )

        def _call() -> str:
            try:
                # Try to force JSON output for better parsing. If unsupported, fall back.
                try:
                    resp = model.generate_content(
                        prompt,
                        generation_config={
                            "temperature": 0.2,
                            "response_mime_type": "application/json",
                        },
                    )
                except TypeError:
                    resp = model.generate_content(prompt)
                return getattr(resp, "text", "") or ""
            except Exception as exc:
                msg = str(exc)
                low = msg.lower()
                if "404" in low or "not found" in low or "listmodels" in low:
                    raise AIServiceError(
                        f"Gemini model topilmadi yoki generateContent ni qo'llamaydi: {model_name}\n"
                        "Yechim: `.env` da `GEMINI_MODEL` ni mavjud modelga almashtiring.\n"
                        "Model ro'yxatini ko'rish: `python scripts\\list_gemini_models.py`"
                    ) from exc
                if _looks_like_gemini_auth_error(msg):
                    raise AIServiceError(_format_gemini_auth_error(msg)) from exc
                raise

        max_retries = int(os.getenv("GEMINI_RETRY_MAX", "1") or 1)
        max_retries = max(0, min(3, max_retries))
        max_delay = int(os.getenv("GEMINI_RETRY_MAX_DELAY", "60") or 60)
        max_delay = max(1, min(120, max_delay))

        raw = ""
        for attempt in range(max_retries + 1):
            try:
                raw = await asyncio.wait_for(asyncio.to_thread(_call), timeout=total_timeout)
                break
            except AIServiceError:
                raise
            except asyncio.TimeoutError as exc:
                if attempt < max_retries:
                    await asyncio.sleep(2 + attempt)
                    continue
                raise AIServiceError(
                    _format_deadline_error("timeout", provider="gemini", model_name=model_name, timeout_sec=total_timeout)
                ) from exc
            except Exception as exc:
                msg = str(exc)
                if _looks_like_gemini_quota_error(msg):
                    delay = _gemini_retry_after_seconds(msg)
                    is_daily = _looks_like_gemini_daily_quota(msg)
                    if (attempt < max_retries) and (not is_daily) and delay and delay <= max_delay:
                        await asyncio.sleep(max(1, int(delay)) + 1)
                        continue
                    raise AIServiceError(_format_gemini_quota_error(msg, model_name=model_name)) from exc
                if _looks_like_deadline_exceeded(msg):
                    if attempt < max_retries:
                        await asyncio.sleep(2 + attempt)
                        continue
                    raise AIServiceError(
                        _format_deadline_error(msg, provider="gemini", model_name=model_name, timeout_sec=total_timeout)
                    ) from exc
                raise
        data = _load_json_from_text(raw)
        return _normalize_quiz(data)

    async def _generate_gemini_topic(
        self,
        topic: str,
        question_count: int,
        *,
        lang_instruction: str,
        extra_instructions: str = "",
    ) -> List[Dict[str, Any]]:
        try:
            import google.generativeai as genai
        except ImportError as exc:
            raise AIServiceError(
                "google-generativeai not installed. Install: pip install google-generativeai"
            ) from exc

        if not self.gemini_api_key:
            raise AIServiceError("GEMINI_API_KEY is not set")

        genai.configure(api_key=self.gemini_api_key)
        model_name = (self.gemini_model or "").strip()
        if not model_name:
            raise AIServiceError("GEMINI_MODEL bo'sh. .env da GEMINI_MODEL=... qo'ying.")
        if model_name.startswith("models/"):
            model_name = model_name.split("/", 1)[1]

        model = genai.GenerativeModel(model_name)

        timeout_sec = int(os.getenv("GEMINI_TIMEOUT_SEC", os.getenv("AI_REQUEST_TIMEOUT_SEC", "60")) or 60)
        timeout_sec = max(5, min(300, timeout_sec))
        total_timeout = int(os.getenv("GEMINI_TOTAL_TIMEOUT_SEC", str(timeout_sec + 30)) or (timeout_sec + 30))
        total_timeout = max(timeout_sec, min(600, total_timeout))

        prompt = (
            "Sen professional testologsan. Berilgan MAVZU bo'yicha ko'p tanlovli test yarat.\n"
            f"- Savollar soni: {question_count}\n"
            "- Har savolda 4 ta variant bo'lsin\n"
            "- Variantlar bir xil uslubda bo'lsin (hammasi ibora yoki hammasi 1 gap)\n"
            "- Variantlar uzunligi bir-biriga yaqin bo'lsin: eng uzun va eng qisqa variant farqi 2-3 so'zdan oshmasin\n"
            "- Juda qisqa (1-2 so'zli) va juda uzun variantlardan qoching; har bir variant taxminan 4-10 so'z bo'lsin\n"
            "- Variantlarda izoh/tushuntirish bo'lmasin; izoh faqat explanation ga yozilsin\n"
            "- Noto'g'ri variantlar ham mantiqan ishonarli bo'lsin, to'g'ri javob ko'zga tashlanib qolmasin\n"
            "- correct_index 0..3 bo'lsin\n"
            "- explanation qisqa bo'lsin\n"
            ""
            f"- Til: {lang_instruction}\n"
            "Natijani faqat JSON ko'rinishida qaytar: {\"quiz\": [...]}.\n"
            "Har element: {\"question\": \"...\", \"options\": [\"A\",\"B\",\"C\",\"D\"], "
            "\"correct_index\": 0, \"explanation\": \"...\"}\n\n"
            f"{extra_instructions}\n\nMavzu: {topic}".strip()
        )

        def _call() -> str:
            try:
                try:
                    resp = model.generate_content(
                        prompt,
                        generation_config={
                            "temperature": 0.2,
                            "response_mime_type": "application/json",
                        },
                    )
                except TypeError:
                    resp = model.generate_content(prompt)
                return getattr(resp, "text", "") or ""
            except Exception as exc:
                msg = str(exc)
                low = msg.lower()
                if "404" in low or "not found" in low or "listmodels" in low:
                    raise AIServiceError(
                        f"Gemini model topilmadi yoki generateContent ni qo'llamaydi: {model_name}\n"
                        "Yechim: `.env` da `GEMINI_MODEL` ni mavjud modelga almashtiring.\n"
                        "Model ro'yxatini ko'rish: `python scripts\\list_gemini_models.py`"
                    ) from exc
                if _looks_like_gemini_auth_error(msg):
                    raise AIServiceError(_format_gemini_auth_error(msg)) from exc
                raise

        max_retries = int(os.getenv("GEMINI_RETRY_MAX", "1") or 1)
        max_retries = max(0, min(3, max_retries))
        max_delay = int(os.getenv("GEMINI_RETRY_MAX_DELAY", "60") or 60)
        max_delay = max(1, min(120, max_delay))

        raw = ""
        for attempt in range(max_retries + 1):
            try:
                raw = await asyncio.wait_for(asyncio.to_thread(_call), timeout=total_timeout)
                break
            except AIServiceError:
                raise
            except asyncio.TimeoutError as exc:
                if attempt < max_retries:
                    await asyncio.sleep(2 + attempt)
                    continue
                raise AIServiceError(
                    _format_deadline_error("timeout", provider="gemini", model_name=model_name, timeout_sec=total_timeout)
                ) from exc
            except Exception as exc:
                msg = str(exc)
                if _looks_like_gemini_quota_error(msg):
                    delay = _gemini_retry_after_seconds(msg)
                    is_daily = _looks_like_gemini_daily_quota(msg)
                    if (attempt < max_retries) and (not is_daily) and delay and delay <= max_delay:
                        await asyncio.sleep(max(1, int(delay)) + 1)
                        continue
                    raise AIServiceError(_format_gemini_quota_error(msg, model_name=model_name)) from exc
                if _looks_like_deadline_exceeded(msg):
                    if attempt < max_retries:
                        await asyncio.sleep(2 + attempt)
                        continue
                    raise AIServiceError(
                        _format_deadline_error(msg, provider="gemini", model_name=model_name, timeout_sec=total_timeout)
                    ) from exc
                raise
        data = _load_json_from_text(raw)
        return _normalize_quiz(data)

    async def _generate_gemini_image(self, image_path: str, *, lang_instruction: str) -> Dict[str, Any]:
        try:
            import google.generativeai as genai
        except ImportError as exc:
            raise AIServiceError(
                "google-generativeai not installed. Install: pip install google-generativeai"
            ) from exc

        if not self.gemini_api_key:
            raise AIServiceError("GEMINI_API_KEY is not set")

        genai.configure(api_key=self.gemini_api_key)
        model_name = (self.gemini_model or "").strip()
        if not model_name:
            raise AIServiceError("GEMINI_MODEL bo'sh. .env da GEMINI_MODEL=... qo'ying.")
        if model_name.startswith("models/"):
            model_name = model_name.split("/", 1)[1]

        model = genai.GenerativeModel(model_name)

        timeout_sec = int(os.getenv("GEMINI_TIMEOUT_SEC", os.getenv("AI_REQUEST_TIMEOUT_SEC", "60")) or 60)
        timeout_sec = max(5, min(300, timeout_sec))
        total_timeout = int(os.getenv("GEMINI_TOTAL_TIMEOUT_SEC", str(timeout_sec + 30)) or (timeout_sec + 30))
        total_timeout = max(timeout_sec, min(600, total_timeout))

        prompt = (
            "Sen professional testologsan. Quyidagi rasmda savol (va ehtimol variantlar) bor.\n"
            "- Agar rasmda 4 ta variant bo'lsa, ularni aynan 4 ta qilib qaytar.\n"
            "- Agar variantlar yo'q bo'lsa, 4 ta mantiqiy variant yarat.\n"
            "- Variantlar bir xil uslubda bo'lsin (hammasi ibora yoki hammasi 1 gap)\n"
            "- Variantlar uzunligi bir-biriga yaqin bo'lsin: eng uzun va eng qisqa variant farqi 2-3 so'zdan oshmasin\n"
            "- Juda qisqa (1-2 so'zli) va juda uzun variantlardan qoching; har bir variant taxminan 4-10 so'z bo'lsin\n"
            "- Variantlarda izoh/tushuntirish bo'lmasin; izoh faqat explanation ga yozilsin\n"
            "- Noto'g'ri variantlar ham mantiqan ishonarli bo'lsin, to'g'ri javob ko'zga tashlanib qolmasin\n"
            "- correct_index 0..3 bo'lsin\n"
            "- explanation qisqa bo'lsin\n"
            f"- Til: {lang_instruction}\n"
            "Natijani faqat JSON ko'rinishida qaytar: {\"quiz\": [...]}.\n"
            "Har element: {\"question\": \"...\", \"options\": [\"A\",\"B\",\"C\",\"D\"], "
            "\"correct_index\": 0, \"explanation\": \"...\"}\n"
            "Eslatma: Agar rasmda to'g'ri javob aniq ko'rinmasa, eng mantiqiy javobni tanla."
        )

        def _call() -> str:
            uploaded = None
            try:
                uploaded = genai.upload_file(image_path)
                try:
                    resp = model.generate_content(
                        [prompt, uploaded],
                        generation_config={
                            "temperature": 0.2,
                            "response_mime_type": "application/json",
                        },
                    )
                except TypeError:
                    resp = model.generate_content([prompt, uploaded])
                return getattr(resp, "text", "") or ""
            finally:
                if uploaded is not None:
                    try:
                        genai.delete_file(uploaded)
                    except Exception:
                        pass

        max_retries = int(os.getenv("GEMINI_RETRY_MAX", "1") or 1)
        max_retries = max(0, min(3, max_retries))
        max_delay = int(os.getenv("GEMINI_RETRY_MAX_DELAY", "60") or 60)
        max_delay = max(1, min(120, max_delay))

        raw = ""
        for attempt in range(max_retries + 1):
            try:
                raw = await asyncio.wait_for(asyncio.to_thread(_call), timeout=total_timeout)
                break
            except asyncio.TimeoutError as exc:
                if attempt < max_retries:
                    await asyncio.sleep(2 + attempt)
                    continue
                raise AIServiceError(
                    _format_deadline_error("timeout", provider="gemini", model_name=model_name, timeout_sec=total_timeout)
                ) from exc
            except Exception as exc:
                msg = str(exc)
                low = msg.lower()
                if "404" in low or "not found" in low or "listmodels" in low:
                    raise AIServiceError(
                        f"Gemini model topilmadi yoki generateContent ni qo'llamaydi: {model_name}\n"
                        "Yechim: `.env` da `GEMINI_MODEL` ni mavjud modelga almashtiring.\n"
                        "Model ro'yxatini ko'rish: `python scripts\\list_gemini_models.py`"
                    ) from exc
                if _looks_like_gemini_auth_error(msg):
                    raise AIServiceError(_format_gemini_auth_error(msg)) from exc
                if _looks_like_gemini_quota_error(msg):
                    delay = _gemini_retry_after_seconds(msg)
                    is_daily = _looks_like_gemini_daily_quota(msg)
                    if (attempt < max_retries) and (not is_daily) and delay and delay <= max_delay:
                        await asyncio.sleep(max(1, int(delay)) + 1)
                        continue
                    raise AIServiceError(_format_gemini_quota_error(msg, model_name=model_name)) from exc
                if _looks_like_deadline_exceeded(msg):
                    if attempt < max_retries:
                        await asyncio.sleep(2 + attempt)
                        continue
                    raise AIServiceError(
                        _format_deadline_error(msg, provider="gemini", model_name=model_name, timeout_sec=total_timeout)
                    ) from exc
                raise

        data = _load_json_from_text(raw)
        questions = _normalize_quiz(data)
        if not questions:
            raise AIServiceError("AI did not return any valid questions from image")
        return questions[0]

